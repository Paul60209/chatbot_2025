import unittest
import asyncio
import os
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tools.translator import translate_text, translate_shape, translate_ppt
from dotenv import load_dotenv
import tempfile

# 載入環境變量
load_dotenv()

class TestTranslator(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        """在所有測試開始前設置環境"""
        if not os.getenv('OPENAI_API_KEY'):
            raise Exception("請設置 OPENAI_API_KEY 環境變量")

    def setUp(self):
        """設置測試環境"""
        # 創建臨時目錄
        self.test_dir = "tests/temp"
        os.makedirs(self.test_dir, exist_ok=True)
        
        # 創建輸出目錄
        self.output_dir = "output"
        os.makedirs(self.output_dir, exist_ok=True)
        
        # 複製測試用 PowerPoint 文件
        self.test_pptx = "tests/test.pptx"
        self.temp_pptx = os.path.join(self.test_dir, "test_copy.pptx")
        if os.path.exists(self.test_pptx):
            shutil.copy2(self.test_pptx, self.temp_pptx)

    def tearDown(self):
        """清理測試環境"""
        # 清理臨時目錄
        if os.path.exists(self.test_dir):
            shutil.rmtree(self.test_dir)

    def test_translate_text(self):
        """測試文本翻譯功能"""
        # 準備測試數據
        text = "你好，世界"
        olang = "zh-TW"
        tlang = "en"
        
        # 執行翻譯
        result = asyncio.run(translate_text(text, olang, tlang))
        
        # 驗證結果
        self.assertIsNotNone(result)
        self.assertIsInstance(result, str)
        self.assertNotEqual(result, text)
        print(f"\n原文: {text}")
        print(f"譯文: {result}")

    def test_translate_shape(self):
        """測試形狀翻譯功能"""
        if not os.path.exists(self.temp_pptx):
            self.skipTest("測試用 PowerPoint 文件不存在")
            
        # 載入測試用 PowerPoint
        prs = Presentation(self.temp_pptx)
        
        # 尋找包含文本的形狀
        text_shape = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    text_shape = shape
                    break
            if text_shape:
                break
                
        if not text_shape:
            self.skipTest("未找到包含文本的形狀")
            
        # 記錄原始文本
        original_text = text_shape.text_frame.text
        print(f"\n原始形狀文本: {original_text}")
        
        # 執行翻譯
        asyncio.run(translate_shape(text_shape, "zh-TW", "en"))
        
        # 驗證結果
        translated_text = text_shape.text_frame.text
        self.assertNotEqual(translated_text, original_text)
        print(f"翻譯後文本: {translated_text}")

    def test_translate_ppt(self):
        """測試 PowerPoint 文件翻譯功能"""
        if not os.path.exists(self.temp_pptx):
            self.skipTest("測試用 PowerPoint 文件不存在")
            
        # 執行翻譯
        output_path = asyncio.run(translate_ppt(self.temp_pptx, "zh-TW", "en"))
        
        # 驗證輸出文件存在
        self.assertTrue(os.path.exists(output_path))
        print(f"\n翻譯後的文件已保存至: {output_path}")
        
        # 載入原始文件和翻譯後的文件
        original_prs = Presentation(self.test_pptx)
        translated_prs = Presentation(output_path)
        
        # 驗證投影片數量相同
        self.assertEqual(len(original_prs.slides), len(translated_prs.slides))
        
        # 驗證至少有一個形狀的文本被翻譯
        found_translation = False
        for slide in translated_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    found_translation = True
                    print(f"找到翻譯後的文本: {shape.text_frame.text}")
                    break
            if found_translation:
                break
                
        self.assertTrue(found_translation, "未找到任何翻譯後的文本")

if __name__ == '__main__':
    unittest.main() 