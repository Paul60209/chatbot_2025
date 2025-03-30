import nest_asyncio
import os
import mimetypes

from pydantic import BaseModel, Field
from typing import Optional, Type
from tools.translator import upload_file, translate_ppt
from langchain.tools import BaseTool
import chainlit as cl
import config
from pptx import Presentation

nest_asyncio.apply()
DELIMITER = os.environ.get('DELIMITER', '||TRANSLATE_DELIMITER||')
open_ai_key = os.getenv('OPENAI_API_KEY', None)

def is_valid_powerpoint(file_path: str) -> bool:
    """檢查文件是否為有效的 PowerPoint 文件。

    Args:
        file_path (str): 文件路徑

    Returns:
        bool: 如果是有效的 PowerPoint 文件則返回 True
    """
    if not file_path:
        return False
        
    # 獲取文件擴展名（轉換為小寫以進行比較）
    ext = os.path.splitext(file_path)[1].lower()
    valid_extensions = ['.ppt', '.pptx']
    
    # 檢查文件擴展名
    if ext not in valid_extensions:
        return False
        
    # 檢查 MIME 類型
    mime_type = mimetypes.guess_type(file_path)[0]
    valid_mime_types = [
        'application/vnd.ms-powerpoint',                    # .ppt
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'  # .pptx
    ]
    
    return mime_type in valid_mime_types

class PowerPointCheckInput(BaseModel):
    """Input for power point translate check."""

    olang: str = Field(...,
                       description="""The original language of the content.""")
    tlang: str = Field(...,
                       description="""The target language for translation.""")


class PowerPointTranslator(BaseTool):
    name = "translate_ppt"
    description = """
            This tool translates PowerPoint files from one language to another.
            Use this tool when the user wants to translate a PowerPoint file and has specified both languages.
            
            Required arguments:
            - olang: The original language code (e.g., 'zh-TW' for Chinese, 'en' for English, 'ja' for Japanese)
            - tlang: The target language code (e.g., 'zh-TW' for Chinese, 'en' for English, 'ja' for Japanese)
            
            The tool will handle the file upload process automatically.
            After translation, it will provide a download link for the translated file.
            
            Supported file formats:
            - .ppt (PowerPoint 97-2003)
            - .pptx (PowerPoint 2007+)
            """
    args_schema: Type[BaseModel] = PowerPointCheckInput

    def _run(self, olang: str, tlang: str):
        """Translates a PowerPoint file.

        Args:
            olang (str): The original language code (e.g., 'zh-TW', 'en', 'ja')
            tlang (str): The target language code (e.g., 'zh-TW', 'en', 'ja')
        """
        import asyncio
        return asyncio.run(self._arun(olang=olang, tlang=tlang))

    async def _arun(self, olang: str, tlang: str):
        """Translates a PowerPoint file asynchronously.

        Args:
            olang (str): The original language code (e.g., 'zh-TW', 'en', 'ja')
            tlang (str): The target language code (e.g., 'zh-TW', 'en', 'ja')
        """
        try:
            # 1. 等待用戶上傳文件
            file_path = await upload_file()
            if not file_path:
                return "請上傳 PowerPoint 文件"
                
            # 2. 驗證文件格式
            if not is_valid_powerpoint(file_path):
                return "請上傳有效的 PowerPoint 文件（.ppt 或 .pptx 格式）"
                
            # 3. 開始翻譯流程
            await cl.Message(content=f"開始將文件從 {olang} 翻譯成 {tlang}...").send()
            
            # 4. 調用翻譯功能
            try:
                await self.translate_ppt(file_path, olang, tlang)
                return "翻譯完成"
            except Exception as e:
                print(f"翻譯過程中發生錯誤: {str(e)}")
                return "翻譯過程中發生錯誤，請稍後再試"
            
        except Exception as e:
            print(f"處理過程中發生錯誤: {str(e)}")
            return f"發生錯誤：{str(e)}"

    async def translate_ppt(self, file_path: str, olang: str, tlang: str) -> str:
        try:
            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)
            
            for i, slide in enumerate(presentation.slides):
                print(f"\n正在翻譯第 {i+1}/{total_slides} 張投影片...")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        original_text = shape.text.strip()
                        print(f"\n原文: {original_text}")
                        
                        run = await self.llm_chain.arun(
                            text=original_text,
                            olang=olang,
                            tlang=tlang
                        )
                        
                        print(f"譯文: {run}")
                        shape.text = run
            return "翻譯完成"
        except Exception as e:
            print(f"翻譯過程中發生錯誤: {str(e)}")
            return "翻譯過程中發生錯誤，請稍後再試"
