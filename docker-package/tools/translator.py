from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_COLOR_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt
import asyncio
import nest_asyncio
import tempfile
import os
import chainlit as cl
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from copy import deepcopy
from langchain.tools import BaseTool
from pydantic import BaseModel, Field
from typing import Type, Optional, Dict, Any
import time

# 定義輸出路徑
OUTPUT_PATH = 'output'

nest_asyncio.apply()

class PowerPointTranslatorInput(BaseModel):
    """PowerPoint 翻譯工具的輸入模型"""
    olang: str = Field(..., description="Original language code (e.g., 'zh-TW', 'en', 'ja')")
    tlang: str = Field(..., description="Target language code (e.g., 'zh-TW', 'en', 'ja')")

class PowerPointTranslator(BaseTool):
    """用於翻譯 PowerPoint 文件的 Langchain 工具"""
    name: str = "translate_ppt"
    description: str = """Translate a PowerPoint file from one language to another.
    
    Args:
        olang (str): The source language code. Must be one of:
            - 'zh-TW' for Chinese
            - 'en' for English
            - 'ja' for Japanese
        tlang (str): The target language code. Must be one of:
            - 'zh-TW' for Chinese
            - 'en' for English
            - 'ja' for Japanese
    
    Returns:
        str: A message indicating the translation status and file location.
    
    Example:
        To translate from Chinese to English:
        translate_ppt(olang="zh-TW", tlang="en")
    """
    args_schema: Type[BaseModel] = PowerPointTranslatorInput

    def _run(self, olang: str, tlang: str) -> str:
        """同步運行方法"""
        loop = asyncio.get_event_loop()
        return loop.run_until_complete(self._arun(olang=olang, tlang=tlang))

    async def _arun(self, olang: str, tlang: str) -> str:
        """異步運行方法，處理 PowerPoint 翻譯請求"""
        try:
            print(f"\nStarting translation tool...")
            print(f"Source language: {olang}")
            print(f"Target language: {tlang}")
            
            # 等待用戶上傳文件
            print("Waiting for file upload...")
            file_path = await upload_file()
            print(f"Upload result: {file_path}")
            
            if not file_path:
                return "No file received or upload failed"

            # 執行翻譯
            print("Starting translation...")
            output_path = await translate_ppt(file_path, olang, tlang)
            print(f"Translation result: {output_path}")
            
            # 返回結果
            if output_path:
                # 發送完成訊息
                await cl.Message(content="翻譯已完成！檔案已儲存至: " + output_path).send()
                return "TRANSLATION_COMPLETE"
            else:
                return "Error occurred during translation"
                
        except Exception as e:
            print(f"Translation tool execution error: {str(e)}")
            return f"Error during translation: {str(e)}"

async def translate_text(text: str, olang: str, tlang: str) -> str:
    """使用 ChatGPT 翻譯文本。

    Args:
        text (str): 要翻譯的文本
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼

    Returns:
        str: 翻譯後的文本
    """
    if not text.strip():
        return text

    print(f"\n正在翻譯文本:")
    print(f"原文 ({olang}): {text}")

    # 創建 ChatGPT 模型
    model = ChatOpenAI(temperature=0)
    
    # 創建系統提示
    system_message = f"""You are a professional translator. Translate the following text from {olang} to {tlang}.
    Rules:
    1. Keep all formatting symbols (like bullet points, numbers) unchanged
    2. Keep all special characters unchanged
    3. Keep all whitespace and line breaks
    4. Only translate the actual text content
    5. Maintain the same tone and style
    6. Do not add any explanations or notes
    7. Keep all numbers and dates unchanged
    8. Keep all proper nouns unchanged unless they have standard translations
    """
    
    # 創建消息列表
    messages = [
        {"role": "system", "content": system_message},
        {"role": "user", "content": text}
    ]
    
    # 執行翻譯
    response = await model.ainvoke(messages)
    translated_text = response.content.strip()
    
    print(f"譯文 ({tlang}): {translated_text}\n")
    return translated_text

def get_text_frame_properties(text_frame):
    """獲取文本框的所有格式屬性"""
    properties = {
        'margin_left': text_frame.margin_left,
        'margin_right': text_frame.margin_right,
        'margin_top': text_frame.margin_top,
        'margin_bottom': text_frame.margin_bottom,
        'vertical_anchor': text_frame.vertical_anchor,
        'word_wrap': text_frame.word_wrap,
        'auto_size': text_frame.auto_size,
    }
    return properties

def get_paragraph_properties(paragraph):
    """獲取段落的所有格式屬性"""
    properties = {
        'alignment': paragraph.alignment,
        'level': paragraph.level,
        'line_spacing': paragraph.line_spacing,
        'space_before': paragraph.space_before,
        'space_after': paragraph.space_after,
    }
    return properties

def get_color_properties(color):
    """獲取顏色屬性"""
    if not color:
        return None
        
    properties = {
        'type': color.type if hasattr(color, 'type') else None,
        'rgb': color.rgb if hasattr(color, 'rgb') else None,
        'theme_color': color.theme_color if hasattr(color, 'theme_color') else None,
        'brightness': color.brightness if hasattr(color, 'brightness') else None,
    }
    return properties

def get_run_properties(run):
    """獲取文本運行的所有格式屬性"""
    font = run.font
    properties = {
        'size': font.size,
        'name': font.name,
        'bold': font.bold,
        'italic': font.italic,
        'underline': font.underline,
        'color': get_color_properties(font.color),
        'fill': get_color_properties(font.fill.fore_color) if hasattr(font, 'fill') else None,
    }
    return properties

def apply_color_properties(color_obj, properties):
    """應用顏色屬性"""
    if not properties or not color_obj:
        return
        
    try:
        # 如果有 RGB 值，直接設置 RGB 顏色
        if properties['rgb']:
            if isinstance(properties['rgb'], (tuple, list)) and len(properties['rgb']) == 3:
                color_obj.rgb = RGBColor(*properties['rgb'])
            else:
                color_obj.rgb = properties['rgb']
        # 如果有主題顏色，設置主題顏色
        elif properties['theme_color'] and properties['theme_color'] != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
            color_obj.theme_color = properties['theme_color']
            if properties['brightness'] is not None:
                color_obj.brightness = properties['brightness']
    except Exception as e:
        print(f"設置顏色時發生錯誤: {str(e)}")
        pass  # 如果設置失敗，保持原有顏色

def apply_text_frame_properties(text_frame, properties):
    """應用文本框格式屬性"""
    text_frame.margin_left = properties['margin_left']
    text_frame.margin_right = properties['margin_right']
    text_frame.margin_top = properties['margin_top']
    text_frame.margin_bottom = properties['margin_bottom']
    text_frame.vertical_anchor = properties['vertical_anchor']
    text_frame.word_wrap = properties['word_wrap']
    text_frame.auto_size = properties['auto_size']

def apply_paragraph_properties(paragraph, properties):
    """應用段落格式屬性"""
    paragraph.alignment = properties['alignment']
    paragraph.level = properties['level']
    paragraph.line_spacing = properties['line_spacing']
    paragraph.space_before = properties['space_before']
    paragraph.space_after = properties['space_after']

def apply_run_properties(run, properties):
    """應用文本運行格式屬性"""
    font = run.font
    if properties['size']:
        font.size = properties['size']
    if properties['name']:
        font.name = properties['name']
    if properties['bold'] is not None:
        font.bold = properties['bold']
    if properties['italic'] is not None:
        font.italic = properties['italic']
    if properties['underline'] is not None:
        font.underline = properties['underline']
    
    # 應用顏色
    if properties['color']:
        apply_color_properties(font.color, properties['color'])
    if properties['fill'] and hasattr(font, 'fill'):
        apply_color_properties(font.fill.fore_color, properties['fill'])

async def translate_group_shape(shape, olang: str, tlang: str) -> None:
    """翻譯群組中的所有形狀。

    Args:
        shape: PowerPoint 群組形狀對象
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
    """
    try:
        if not hasattr(shape, 'shapes'):
            return
            
        # 遍歷群組中的所有形狀
        for child_shape in shape.shapes:
            if child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 遞歸處理嵌套群組
                await translate_group_shape(child_shape, olang, tlang)
            else:
                # 翻譯單個形狀
                await translate_shape(child_shape, olang, tlang)
    except Exception as e:
        print(f"翻譯群組形狀時發生錯誤: {str(e)}")
        raise

async def translate_shape(shape, olang: str, tlang: str) -> None:
    """翻譯 PowerPoint 中的形狀。

    Args:
        shape: PowerPoint 形狀對象
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼
    """
    try:
        # 處理群組形狀
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            await translate_group_shape(shape, olang, tlang)
            return
            
        # 檢查形狀是否包含文本框
        if not hasattr(shape, "text_frame"):
            return
            
        text_frame = shape.text_frame
        if not text_frame.text.strip():
            return
            
        # 保存文本框格式
        text_frame_props = get_text_frame_properties(text_frame)
        
        # 遍歷所有段落
        for paragraph in text_frame.paragraphs:
            # 保存段落格式
            para_props = get_paragraph_properties(paragraph)
            
            # 遍歷所有文本運行
            runs_data = []
            for run in paragraph.runs:
                # 保存運行格式和文本
                run_props = get_run_properties(run)
                original_text = run.text
                if original_text.strip():
                    translated_text = await translate_text(original_text, olang, tlang)
                    runs_data.append((translated_text, run_props))
                else:
                    runs_data.append((original_text, run_props))
            
            # 清除原有內容
            for _ in range(len(paragraph.runs)):
                paragraph._p.remove(paragraph.runs[0]._r)
            
            # 添加翻譯後的文本並應用格式
            for text, props in runs_data:
                run = paragraph.add_run()
                run.text = text
                apply_run_properties(run, props)
            
            # 恢復段落格式
            apply_paragraph_properties(paragraph, para_props)
        
        # 恢復文本框格式
        apply_text_frame_properties(text_frame, text_frame_props)
        
    except Exception as e:
        print(f"翻譯形狀時發生錯誤: {str(e)}")
        raise

async def translate_ppt(file_path: str, olang: str, tlang: str) -> str:
    """翻譯 PowerPoint 文件。

    Args:
        file_path (str): PowerPoint 文件路徑
        olang (str): 原始語言代碼
        tlang (str): 目標語言代碼

    Returns:
        str: 翻譯後的文件路徑
    """
    try:
        # 1. 建立輸出目錄
        os.makedirs(OUTPUT_PATH, exist_ok=True)
        
        # 2. 準備輸出文件路徑
        file_name = os.path.basename(file_path)
        name, ext = os.path.splitext(file_name)
        output_file = f'translated_{name}{ext}'
        output_path = os.path.join(OUTPUT_PATH, output_file)
        
        # 3. 載入 PowerPoint
        print("\nStarting PowerPoint translation...")
        print(f"Source language: {olang}")
        print(f"Target language: {tlang}")
        await cl.Message(content=f"Starting translation...\nFrom {olang} to {tlang}").send()
        
        presentation = Presentation(file_path)
        total_slides = len(presentation.slides)
        
        # 4. 翻譯每個投影片
        for index, slide in enumerate(presentation.slides, 1):
            progress_msg = f"Translating slide {index}/{total_slides}..."
            print(f"\n{progress_msg}")
            await cl.Message(content=progress_msg).send()
            for shape in slide.shapes:
                await translate_shape(shape, olang, tlang)
        
        # 5. 儲存翻譯後的文件
        print("\nSaving translated file...")
        await cl.Message(content="Translation completed, generating file...").send()
        presentation.save(output_path)
        
        # 6. 刪除臨時文件
        if os.path.exists(file_path):
            os.remove(file_path)
        
        try:
            # 7. 建立下載連結（僅在 Chainlit 環境中）
            elements = [
                cl.File(
                    name=output_file,
                    path=output_path,
                    display="inline"
                )
            ]
            
            # 8. 發送完成消息和下載連結
            await cl.Message(
                content="Translation completed! Click the link below to download the translated file:",
                elements=elements
            ).send()
        except Exception as e:
            # 在非 Chainlit 環境中，只打印消息
            print("\nTranslation completed! File saved to:", output_path)
        
        return output_path
        
    except Exception as e:
        print(f"\nError during translation process: {str(e)}")
        if os.path.exists(file_path):
            os.remove(file_path)
        raise

async def upload_file() -> str:
    """處理文件上傳。

    Returns:
        str: 上傳的文件路徑
    """
    try:
        # 等待用戶上傳文件
        files = await cl.AskFileMessage(
            content="Please upload a PowerPoint file (.ppt or .pptx, max file size: 10MB)",
            accept=["application/vnd.ms-powerpoint", 
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation"],
            max_size_mb=10,
            timeout=20
        ).send()

        if not files or len(files) == 0:
            print("No file received")
            return None

        file = files[0]
        print(f"\nFile received: {file.name}")
        print(f"File object attributes: {dir(file)}")  # 打印文件物件的所有屬性
        
        # 檢查檔案副檔名
        file_name = file.name.lower()
        if not file_name.endswith(('.ppt', '.pptx')):
            print(f"Unsupported file format: {file_name}")
            await cl.Message(content="Please upload a .ppt or .pptx file").send()
            return None

        # 建立臨時檔案
        temp_dir = tempfile.gettempdir()
        temp_file_path = os.path.join(temp_dir, file.name)
        
        # 寫入檔案內容
        with open(temp_file_path, 'wb') as f:
            if hasattr(file, 'path'):
                # 如果檔案已經在本地，直接複製
                with open(file.path, 'rb') as source:
                    f.write(source.read())
            elif hasattr(file, 'content'):
                print(f" use File content")
                f.write(file.content)
            elif hasattr(file, 'bytes'):
                print(f" use File bytes")
                f.write(file.bytes)
            else:
                print("No valid file content found")
                return None
            
        print(f"File saved to: {temp_file_path}")
        return temp_file_path
        
    except Exception as e:
        print(f"Error during file upload: {str(e)}")
        return None 